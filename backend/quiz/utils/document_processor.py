"""Document processing utilities"""
import os
import pytesseract
from PIL import Image
import PyPDF2
from pptx import Presentation
from docx import Document as DocxDocument
from io import BytesIO


class DocumentProcessor:
    """Process various document types and extract text"""

    @staticmethod
    def extract_text_from_pdf(file):
        """Extract text from PDF"""
        text = ""
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()

    @staticmethod
    def extract_text_from_pptx(file):
        """Extract text from PPTX"""
        text = ""
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()

    @staticmethod
    def extract_text_from_docx(file):
        """Extract text from DOCX"""
        doc = DocxDocument(file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()

    @staticmethod
    def extract_text_from_txt(file):
        """Extract text from TXT"""
        file.seek(0)
        return file.read().decode('utf-8').strip()

    @staticmethod
    def extract_text_from_image(file):
        """Extract text from image using OCR"""
        image = Image.open(file)
        text = pytesseract.image_to_string(image)
        return text.strip()

    @classmethod
    def process_document(cls, file, file_type):
        """Process document based on file type"""
        file.seek(0)
        
        if file_type == 'pdf':
            return cls.extract_text_from_pdf(file)
        elif file_type == 'pptx' or file_type == 'ppt':
            return cls.extract_text_from_pptx(file)
        elif file_type == 'docx' or file_type == 'doc':
            return cls.extract_text_from_docx(file)
        elif file_type == 'txt':
            return cls.extract_text_from_txt(file)
        elif file_type in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
            return cls.extract_text_from_image(file)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
