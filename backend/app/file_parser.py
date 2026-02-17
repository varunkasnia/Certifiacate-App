from io import BytesIO
from pathlib import Path

from fastapi import HTTPException, UploadFile
from PyPDF2 import PdfReader
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".ppt", ".pptx", ".md"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}


async def parse_upload_to_text(file: UploadFile) -> tuple[str, str, bytes]:
    filename = file.filename or "upload"
    ext = Path(filename).suffix.lower()
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    if ext in IMAGE_EXTENSIONS:
        return "", ext, content

    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Unsupported file type. Use PDF/TXT/PPT/PPTX/MD/Images")

    try:
        if ext in {".txt", ".md"}:
            return content.decode("utf-8", errors="ignore"), ext, content
        if ext == ".pdf":
            return _parse_pdf(content), ext, content
        if ext in {".ppt", ".pptx"}:
            return _parse_pptx(content), ext, content
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Could not parse file: {exc}") from exc

    raise HTTPException(status_code=400, detail="Unsupported file type")


def _parse_pdf(content: bytes) -> str:
    reader = PdfReader(BytesIO(content))
    text_parts = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(text_parts).strip()


def _parse_pptx(content: bytes) -> str:
    presentation = Presentation(BytesIO(content))
    chunks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks).strip()
